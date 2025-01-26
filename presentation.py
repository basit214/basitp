from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize presentation
presentation = Presentation()

# Function to add a slide with title and content
def add_slide(presentation, title, content, bullet_points=False):
    slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content_placeholder = slide.placeholders[1]
    if bullet_points:
        content_placeholder.text = ""
        for line in content.split("\n"):
            p = content_placeholder.text_frame.add_paragraph()
            p.text = line
            p.level = 0
    else:
        content_placeholder.text = content

# Function to add a title slide
def add_title_slide(presentation, title, subtitle):
    slide_layout = presentation.slide_layouts[0]  # Title Slide layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Add Title Slide
add_title_slide(
    presentation,
    "AWS Bootcamp: Week 1 - AWS Fundamentals and Core Services",
    "Presented by: [Your Name/Organization]\nDate: [Your Date]",
)

# Week 1: Day 1 Content
add_slide(
    presentation,
    "Introduction to Cloud Computing and AWS",
    "What is Cloud Computing?\n- Benefits of Cloud Computing\n\nIntroduction to AWS:\n- What is it? Why AWS?\n\nAWS Global Infrastructure:\n- Overview and Regions\n\nAWS Free Tier:\n- Overview and Pricing",
    bullet_points=True,
)

# Week 1: Day 2-3 Content
add_slide(
    presentation,
    "Compute Services (EC2, Lambda, Elastic Beanstalk)",
    "EC2:\n- Instances, Types, Pricing, and Management\n- Launching and Managing EC2 Instances\n\nLambda:\n- Serverless architecture and use cases\n\nElastic Beanstalk:\n- Deploying applications without managing infrastructure",
    bullet_points=True,
)

# Week 1: Day 4-5 Content
add_slide(
    presentation,
    "Storage Services (S3, EBS, Glacier, EFS)",
    "AWS Storage Options:\n- S3: Buckets, Objects, Versioning, and Permissions\n- EBS: Persistent block storage for EC2\n- Glacier: Archival storage\n- EFS: File storage and use cases",
    bullet_points=True,
)

# Week 1: Day 6-7 Content
add_slide(
    presentation,
    "Networking Services (VPC, Route 53, ELB, CloudFront)",
    "VPC:\n- Virtual Private Cloud, Subnets, Route Tables, and Security Groups\n\nRoute 53:\n- DNS Management and Domain Setup\n\nElastic Load Balancer (ELB):\n- Load Balancing and Content Delivery\n\nCloudFront:\n- Content Delivery Network (CDN)",
    bullet_points=True,
)

# Week 2: Advanced Topics
add_slide(
    presentation,
    "Advanced Topics and Real-World Applications",
    "Security and Identity Management (IAM, KMS, Secrets Manager):\n- IAM: Users, Groups, Roles, and Policies\n- KMS: Key Management\n- Secrets Manager: Managing secrets and credentials\n\nDatabases (RDS, DynamoDB, Aurora, Redshift):\n- RDS: Relational Databases\n- DynamoDB: NoSQL Database\n- Aurora and Redshift: High-performance Databases",
    bullet_points=True,
)

# Week 2: DevOps Tools
add_slide(
    presentation,
    "DevOps and Automation Tools",
    "Infrastructure as Code:\n- Using AWS CloudFormation\n\nCI/CD Tools:\n- CodePipeline and CodeDeploy\n\nAutomation:\n- Automating deployments and scaling applications",
    bullet_points=True,
)

# Week 2: Final Preparation
add_slide(
    presentation,
    "Job Preparation and AWS Certification",
    "AWS Certifications Overview:\n- Solutions Architect, Developer, SysOps, etc.\n\nResume Building:\n- Highlighting AWS skills\n\nInterview Preparation:\n- Common AWS questions and scenarios\n\nHands-on Practice:\n- Final mini-project or case study",
    bullet_points=True,
)

# Save Presentation
output_path = "AWS_Bootcamp_Presentation.pptx"
presentation.save(output_path)
output_path
